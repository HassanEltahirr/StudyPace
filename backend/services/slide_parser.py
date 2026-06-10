from __future__ import annotations

import io
import re
import zipfile
from dataclasses import dataclass
from pathlib import Path
from xml.etree import ElementTree

try:
    import fitz  # PyMuPDF
except Exception:  # pragma: no cover - optional runtime dependency
    fitz = None

try:
    from pptx import Presentation
except Exception:  # pragma: no cover - optional runtime dependency
    Presentation = None


class SlideExtractionError(ValueError):
    pass


@dataclass
class ExtractedSlide:
    slide_number: int
    title: str
    text: str
    content_tags: list[str]


BULLET_CHARS = "•◦◼■▪▫❑□●○▶►▸"


def extract_slides(filename: str, content: bytes) -> list[ExtractedSlide]:
    ext = Path(filename).suffix.lower()
    if ext == ".pdf":
        return _extract_pdf(content)
    if ext in {".pptx", ".ppt"}:
        return _extract_pptx(content)
    if ext == ".docx":
        return _extract_docx(content)
    if ext in {".txt", ".md"}:
        return _extract_text(content)

    try:
        return _extract_text(content)
    except UnicodeDecodeError as exc:
        raise SlideExtractionError("Unsupported file type. Upload PDF, PPTX, DOCX, TXT, or Markdown.") from exc


def _extract_pdf(content: bytes) -> list[ExtractedSlide]:
    if fitz is None:
        raise SlideExtractionError("PDF extraction needs PyMuPDF. Rebuild the backend after installing requirements.txt.")

    try:
        doc = fitz.open(stream=content, filetype="pdf")
    except Exception as exc:
        raise SlideExtractionError("Could not read this PDF. It may be encrypted or corrupted.") from exc

    slides: list[ExtractedSlide] = []
    for idx, page in enumerate(doc, start=1):
        text = page.get_text("text").strip()
        if text:
            slides.append(_make_slide(idx, text))
        else:
            slides.append(_make_visual_slide(idx))
    return _validate(slides)


def _extract_pptx(content: bytes) -> list[ExtractedSlide]:
    if Presentation is None:
        raise SlideExtractionError("PPTX extraction needs python-pptx. Rebuild the backend after installing requirements.txt.")

    try:
        deck = Presentation(io.BytesIO(content))
    except Exception as exc:
        raise SlideExtractionError("Could not read this PPTX. It may be corrupted or password protected.") from exc

    slides: list[ExtractedSlide] = []
    for idx, slide in enumerate(deck.slides, start=1):
        parts: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                parts.append(shape.text)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if cells:
                        parts.append(" | ".join(cells))
        text = "\n".join(parts).strip()
        if text:
            slides.append(_make_slide(idx, text))
    return _validate(slides)


def _extract_text(content: bytes) -> list[ExtractedSlide]:
    raw = content.decode("utf-8")
    # Explicit slide separators (--- or form feed)
    chunks = [c.strip() for c in re.split(r"\n\s*---+\s*\n|\f", raw) if c.strip()]
    # Markdown headings give natural section boundaries
    if len(chunks) <= 1:
        heading_chunks = [c.strip() for c in re.split(r"(?m)(?=^#{1,3}\s)", raw) if c.strip()]
        if len(heading_chunks) > 1:
            chunks = heading_chunks
    # Character-length fallback for unstructured plain text
    if len(chunks) <= 1 and len(raw) > 1600:
        chunks = [raw[i:i + 1400].strip() for i in range(0, len(raw), 1400) if raw[i:i + 1400].strip()]
    if not chunks:
        chunks = [raw.strip()]
    return _validate([_make_slide(i + 1, chunk) for i, chunk in enumerate(chunks) if chunk.strip()])


def _extract_docx(content: bytes) -> list[ExtractedSlide]:
    try:
        with zipfile.ZipFile(io.BytesIO(content)) as docx:
            document_xml = docx.read("word/document.xml")
    except Exception as exc:
        raise SlideExtractionError("Could not read this DOCX. It may be corrupted or password protected.") from exc

    try:
        root = ElementTree.fromstring(document_xml)
    except ElementTree.ParseError as exc:
        raise SlideExtractionError("Could not parse text from this DOCX.") from exc

    namespace = {"w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main"}
    lines: list[str] = []
    for paragraph in root.findall(".//w:p", namespace):
        text = "".join(node.text or "" for node in paragraph.findall(".//w:t", namespace))
        cleaned = re.sub(r"\s+", " ", text).strip()
        if cleaned:
            lines.append(cleaned)

    raw = "\n".join(lines)
    if not raw.strip():
        raise SlideExtractionError("No usable DOCX text was found. Confirm the syllabus contains selectable text.")
    return _extract_text(raw.encode("utf-8"))


def _make_slide(number: int, text: str) -> ExtractedSlide:
    lines = _clean_lines(text)
    title = _detect_title(lines, number)
    tags = _detect_tags(text)
    return ExtractedSlide(
        slide_number=number,
        title=title,
        text="\n".join(lines),
        content_tags=tags,
    )


def _make_visual_slide(number: int) -> ExtractedSlide:
    return ExtractedSlide(
        slide_number=number,
        title=f"Slide {number}",
        text=f"This PDF slide is image-based. Open the rendered slide image to study slide {number}.",
        content_tags=["visual-only"],
    )


def _clean_lines(text: str) -> list[str]:
    lines = []
    for line in text.replace("\x00", " ").splitlines():
        cleaned = re.sub(r"\s+", " ", line).strip(f" \t-{BULLET_CHARS}")
        cleaned = re.sub(rf"^[{re.escape(BULLET_CHARS)}]+\s*", "", cleaned)
        cleaned = re.sub(r"^#{1,6}\s*", "", cleaned)
        if cleaned:
            lines.append(cleaned)
    return lines


def _detect_title(lines: list[str], number: int) -> str:
    if not lines:
        return f"Slide {number}"
    short_lines = [line for line in lines[:5] if 3 <= len(line) <= 90]
    return short_lines[0] if short_lines else f"Slide {number}"


def _detect_tags(text: str) -> list[str]:
    lower = text.lower()
    tags: list[str] = []
    checks = [
        ("learning objective", ["objective", "outcome", "you will learn"]),
        ("definition", ["define", "definition", " is a ", " refers to "]),
        ("formula", ["=", "∑", "formula", "theorem", "equation"]),
        ("example", ["example", "case study", "for instance"]),
        ("diagram caption", ["figure", "diagram", "table", "chart"]),
        ("question", ["?", "question", "exercise", "problem", "checkpoint", "quiz yourself"]),
        ("numerical problem", ["calculate", "compute", "solve", "find", "determine", "derive", "evaluate", "estimate"]),
    ]
    for tag, needles in checks:
        if any(needle in lower for needle in needles):
            tags.append(tag)
    return tags or ["concept"]


def _validate(slides: list[ExtractedSlide]) -> list[ExtractedSlide]:
    useful = [s for s in slides if len(s.text.split()) >= 5]
    if not useful:
        raise SlideExtractionError("No usable slide text was found. Try a text-based PDF/PPTX or confirm the slide content manually.")
    return useful
